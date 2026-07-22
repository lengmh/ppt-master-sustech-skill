from __future__ import annotations

import subprocess
import sys
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SCRIPTS = Path(__file__).resolve().parents[1]


def write_flat_spec_lock(project: Path) -> None:
    (project / "spec_lock.md").write_text(
        """<!-- ppt-master-schema: spec-lock/v1 -->
# Execution Lock

## canvas
- viewBox: 0 0 960 720
- format: PPT 4:3

## communication
- audience: Regression test
- objective: Verify the v4.1 native round-trip contract.
- core_message: Native replacements and timing survive the round trip.

## mode
- mode: briefing

## visual_style
- visual_style: soft-rounded

## colors
- bg: #FFFFFF
- primary: #2563EB
- accent: #F59E0B
- text: #111827

## typography
- font_family: Aptos
- title_family: Aptos Display
- body_family: Aptos
- body: 18
- title: 28

## icons
- library: phosphor-duotone
- inventory: none

## page_rhythm
- P01: dense

## pptx_structure
- mode: flat

## forbidden
- Mixing icon libraries
""",
        encoding="utf-8",
    )


def run_cli(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, *args],
        check=False,
        capture_output=True,
        text=True,
    )


def test_preset_shape_cli_emits_v41_compact_authored_contract() -> None:
    result = run_cli(
        str(SCRIPTS / "preset_shape_svg.py"),
        "render",
        "rightArrow",
        "--id",
        "process-arrow",
        "--frame",
        "120",
        "180",
        "240",
        "96",
        "--fill",
        "#2563EB",
    )
    assert result.returncode == 0, result.stderr
    root = ET.fromstring(
        f'<svg xmlns="http://www.w3.org/2000/svg">{result.stdout}</svg>'
    )
    group = next(iter(root))
    assert group.get("id") == "process-arrow"
    assert group.get("data-pptx-authoring") == "preset"
    assert group.get("data-pptx-object") == "shape"
    assert group.get("data-pptx-prst") == "rightArrow"
    assert group.get("data-pptx-frame") == "120 180 240 96"
    # v4.1 compact authored presets deliberately omit import transport fields.
    assert group.get("data-pptx-preview-sha256") is None
    assert all(node.get("data-pptx-part") is None for node in group)


def test_pptx_svg_pptx_roundtrip_preserves_v41_replacements_and_timing(
    tmp_path: Path,
) -> None:
    source = tmp_path / "native-source.pptx"
    output = tmp_path / "native-svg"
    roundtrip = tmp_path / "native-roundtrip.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(0.5),
        Inches(0.5),
        Inches(2.0),
        Inches(0.8),
    )
    table = slide.shapes.add_table(
        2,
        2,
        Inches(0.5),
        Inches(1.7),
        Inches(3.0),
        Inches(1.2),
    ).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "10"

    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("Series 1", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4.0),
        Inches(0.8),
        Inches(4.0),
        Inches(2.8),
        chart_data,
    )
    presentation.save(source)

    result = run_cli(
        str(SCRIPTS / "pptx_to_svg.py"),
        str(source),
        "-o",
        str(output),
        "--inheritance-mode",
        "flat",
    )
    assert result.returncode == 0, result.stderr
    svg_files = sorted((output / "svg").glob("*.svg"))
    assert svg_files
    svg = "\n".join(path.read_text(encoding="utf-8") for path in svg_files)
    assert 'data-pptx-prst="rightArrow"' in svg
    assert 'data-pptx-replace-with="table"' in svg
    assert 'data-pptx-replace-with="chart"' in svg
    assert 'data-pptx-fallback-sha256=' in svg

    repeated = run_cli(
        str(SCRIPTS / "pptx_to_svg.py"),
        str(source),
        "-o",
        str(output),
        "--inheritance-mode",
        "flat",
    )
    assert repeated.returncode == 0, repeated.stderr
    write_flat_spec_lock(output)

    exported = run_cli(
        str(SCRIPTS / "svg_to_pptx.py"),
        str(output),
        "-s",
        "svg",
        "-o",
        str(roundtrip),
        "--native-charts-and-tables",
        "--transition",
        "fade",
        "--animation",
        "fade",
        "--animation-trigger",
        "after-previous",
    )
    assert exported.returncode == 0, exported.stderr
    assert len(Presentation(roundtrip).slides) == 1

    with ZipFile(roundtrip) as package:
        names = set(package.namelist())
        slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

    assert '<a:prstGeom prst="rightArrow"' in slide_xml
    assert "<a:tbl>" in slide_xml
    assert "<c:chart" in slide_xml
    assert "<p:transition" in slide_xml
    assert "<p:timing" in slide_xml
    assert any(name.startswith("ppt/charts/chart") for name in names)
    assert any(name.startswith("ppt/embeddings/") for name in names)
